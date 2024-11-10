import numpy as np
import pandas as pd
from datetime import date
from pptx import Presentation   
from pptx.util import Inches, Pt   
from pptx.dml.color import RGBColor  

#### 價格 地段  用途表格
"""
return 會輸出 用途0~4的 價格&地段的表
"""
def detail(df,places,ages):
    unit_place_age0=[]  ## 用途 0
    unit_place_age1=[]  ## 用途 1
    unit_place_age2=[]  ## 用途 2
    unit_place_age3=[]  ## 用途 3
    unit_place_age4=[]  ## 用途 4
    unit_price_use=[]
    # 用來控地點
    for i,place in enumerate (places):
        # age
        for j,age in enumerate(ages):
            # i用來控制用途
            for u in range(5):
                unit_price_use.append(df["unit_price"][(df["地段"]==i) &(df["age"]==age)&(df["主要用途"]==u)].sum()/len(df["unit_price"][(df["地段"]==i) &(df["age"]==age)&(df["主要用途"]==u)]))
        
    unit_price_use = np.array(unit_price_use).reshape(len(places),4,5)
    
    for j , p in enumerate (places):
        unit_place_age0.append({"地點":p,ages[0]:unit_price_use[j][:,0][0], ages[1]:unit_price_use[j][:,0][1],ages[2]:unit_price_use[j][:,0][2],ages[3]:unit_price_use[j][:,0][3]})
        unit_place_age1.append({"地點":p,ages[0]:unit_price_use[j][:,1][0], ages[1]:unit_price_use[j][:,1][1],ages[2]:unit_price_use[j][:,1][2],ages[3]:unit_price_use[j][:,1][3]})
        unit_place_age2.append({"地點":p,ages[0]:unit_price_use[j][:,2][0], ages[1]:unit_price_use[j][:,2][1],ages[2]:unit_price_use[j][:,2][2],ages[3]:unit_price_use[j][:,2][3]})
        unit_place_age3.append({"地點":p,ages[0]:unit_price_use[j][:,3][0], ages[1]:unit_price_use[j][:,3][1],ages[2]:unit_price_use[j][:,3][2],ages[3]:unit_price_use[j][:,3][3]})
        unit_place_age4.append({"地點":p,ages[0]:unit_price_use[j][:,4][0], ages[1]:unit_price_use[j][:,4][1],ages[2]:unit_price_use[j][:,4][2],ages[3]:unit_price_use[j][:,4][3]})

    unit_place_age_df0 = pd.DataFrame(unit_place_age0)
    unit_place_age_df1 = pd.DataFrame(unit_place_age1)
    unit_place_age_df2 = pd.DataFrame(unit_place_age2)
    unit_place_age_df3 = pd.DataFrame(unit_place_age3)
    unit_place_age_df4 = pd.DataFrame(unit_place_age4)
    unit_place_age_df0.set_index('地點', inplace=True)
    unit_place_age_df1.set_index('地點', inplace=True)
    unit_place_age_df2.set_index('地點', inplace=True)
    unit_place_age_df3.set_index('地點', inplace=True)
    unit_place_age_df4.set_index('地點', inplace=True)
    return unit_place_age_df0,unit_place_age_df1,unit_place_age_df2,unit_place_age_df3,unit_place_age_df4

"""
特定給 屋齡 與 用途
"""
def sumX(df,ages,length,usefor):
    scores=[]
    for i in range(4):
        scores.append({"屋齡(年)":ages[i],usefor[0]:length[i][0], usefor[1]:length[i][1],usefor[2]:length[i][2],usefor[3]:length[i][3],usefor[4]:length[i][4]})
        #scores.append({"屋齡":ages[4],usefor[0]:length[l][0], usefor[1]:length[l][1],usefor[2]:length[l][2],usefor[3]:length[l][3]})
    score_df = pd.DataFrame(scores)
    return score_df

"""
特定給 建物的年齡以及用途
"""
def age_use2(df,places,ages,option,usefor):
    P,kind_P,P_age= 0,0,0
    p,kind_p,p_age= float("inf"),0,0
    show = []  
    length,place_price = [],[]
    # age
    for j,age in enumerate(ages):
        # i用來控制用途
        for i in range(5):
            length.append(len(df[(df["地段"]==places[option]) &(df["age"]==age) & (df["主要用途"]==int(i))]))
            place_price.append(df["total_price"][(df["地段"]==places[option]) &(df["age"]==age) & (df["主要用途"]==int(i))].sum())
            
            tmp = df["unit_price"][(df["地段"]==places[option]) &(df["age"]==age) & (df["主要用途"]==int(i))].sum()/len(df["unit_price"][(df["地段"]==places[option]) &(df["age"]==age) & (df["主要用途"]==int(i))])
            price_sum = df["unit_price"][(df["地段"]==int(0)) &(df["age"]==age) & (df["主要用途"]==int(i))].sum()
            number = len(df["unit_price"][(df["地段"]==places[option]) &(df["age"]==age) & (df["主要用途"]==int(i))])

            if P < tmp :
                P = tmp
                kind_P = i
                P_age = age
                P_number = number
            if p > tmp :
                p = tmp
                kind_p = i
                p_age = age
                p_number = number

    length = np.array(length).reshape(4,5)
    place_price = np.array(place_price).reshape(4,5)
    return length,place_price,(P,usefor[kind_P],P_age),(p,usefor[kind_p],p_age)


def aa(a0,ages):
    p=[]
    for i,age in enumerate(ages):
        p.append({"屋齡(年)":age,"地點":str(a0[a0[age] == a0[age].max()].index.values).replace("'",''),"坪/萬":a0[age][a0[age] == a0[age].max()].values })
    p = pd.DataFrame(p)
    p.set_index('屋齡(年)', inplace=True)
    return p

"""
prs = genPPTX(mainTitle,subTitle): 產生一份新的投影片
"""
def genPPTX(mainTitle,subTitle):            
    prs = Presentation()                 
    slide0 = prs.slides.add_slide(prs.slide_layouts[0])               
    slide0.shapes.title.text = mainTitle
    slide0.placeholders[1].text = subTitle
    return prs

"""
prs = addBulletPage(prs,Ptitle,Plist,Plevel): 
    增加一個重點(Plist)頁,並設定重點層級(Plevel)及顏色 (Plevel=1)
"""
def addBulletPage(prs,Ptitle,Plist,Plevel):
    
    slide = prs.slides.add_slide( prs.slide_layouts[1] );  #-- 產生一頁(slide)新的 "標題與內容" 的重點頁(BulletPage)  
    slide.shapes.title.text = Ptitle                       #-- 設定標題(Ptitle)
    tf = slide.shapes.placeholders[1].text_frame           #-- 設定內文 文字框(tf)
    for k in np.arange(len(Plist)):
        if k==0:
            tf.text = Plist[0]   #-- 設定第 1 子標題 (tf.text = Plist[0])
        else:                    #-- 設定新增 子標題 (Plist[k]), 其層級 (Plevel[k]) 及顏色 (Plevel=1為粗體彩色)
            p = tf.add_paragraph()   
            p.level = Plevel[k]
            p.text = Plist[k]   
            if (p.level==1): 
                p.font.bold = True
                p.font.color.rgb = RGBColor(0,0,255)  # RGBColor(0xFF, 0x7F, 0x50)

    print("addBulletPage>>> generate Bullet Page-"+Ptitle)     
    return prs

"""
prs = addSlideDF(prs,ind,Ptable): 
    將表格(Ptable)加入某頁 (prs.slides[ind])
"""
def addSlideDF(prs,ind,Ptable):              
    shapes = prs.slides[ind].shapes

    if (Ptable is not None):
        print("addSlideDF>>> generate dataframe Table...")
        ## location
        left, top, width, height = Inches(1), Inches(1), Inches(8), Inches(6)
        table = shapes.add_table(Ptable.shape[0],
                                 Ptable.shape[1],
                                 left, top, width, height).table
              
        for i in np.arange(Ptable.shape[0]):
            for j in np.arange(Ptable.shape[1]):
                table.cell(i,j).text = str(list(Ptable.iloc[i])[j])
    return prs

def add_personal_DF(prs,width,height,ind,df:pd.DataFrame , use_columns):
    shapes = prs.slides[ind].shapes
    
    rows = df.shape[0]
    cols = len(use_columns)
 
    # table size & location
    left = top = Inches(2.0)
    
 
    # add table into slide
    table = shapes.add_table(rows+1, cols, left, top, width, height).table
    ## table colum
    
    for j in range(len(use_columns)):
        table.cell(0, j).text = str(use_columns[j])
    # table style
    table.style = "GRADE TABLE"
    ## table start
    for i in range(rows):
        for j in range(cols):
            table.cell(i+1, j).text = str(list(df.iloc[i])[j])
    return prs

def makeDFtable(df):                         ##==> table = makeDFtable(df): make df to table with first row as column names
    Xcol = pd.DataFrame(df.columns).transpose()
    Xcol.columns = df.columns
    AAA  = pd.concat([Xcol,df],axis=0)
    Arow = pd.DataFrame(AAA.index)
    Arow.index = Arow[Arow.columns[0]]
    BBB  = pd.concat([Arow,AAA],axis=1)
    BBB.index = Arow[Arow.columns[0]]
    return BBB

"""
prs = 建立好的ppt 還沒save的
"""
def create_ppt(raw,df,data_explain,ppt_name,places,use_columns,translate_raw,ages):
    prs = genPPTX(ppt_name,"by Lawrence,童舟,陳姿諭\n"+str(date.today()))
    
    # p1
    prs = addBulletPage(prs,"KDD1 載入原始數據",
                        ["1-1 原始數據展示"],
                        [0])
    prs = add_personal_DF(prs,Inches(6.0),Inches(0.8), 1, raw,use_columns)
    # -------------------------------------
    # p2
    data_explain = pd.DataFrame(data_explain)
    prs = addBulletPage(prs,"KDD1 載入原始數據",
                        ["1-2 數據說明"],
                        [0])
    prs = addSlideDF(prs, 2, data_explain)
    # -------------------------------------
    # -------------------------------------
    # p3
    prs = addBulletPage(prs,"KDD2 探索交易數據",
                        ["2-1 correlation"],
                        [0])
    # 說明請後製
    prs.slides[3].shapes.add_picture("img/correlation.png", Inches(2), Inches(2))
    # -------------------------------------
    # -------------------------------------
    # p4
    
    prs = addBulletPage(prs,"KDD3 交易數據轉換",
                        ["3-1 產生的數據標籤"],
                        [0])
    translate_raw = translate_raw.head(2)
    translate_raw_cols = np.array(translate_raw.columns)
    # 說明請後製  
    prs = add_personal_DF(prs,Inches(8.0),Inches(0.8), 4, translate_raw,translate_raw_cols)
    # -------------------------------------
    # -------------------------------------
    # p5
    prs = addBulletPage(prs,"KDD4 交易模型（一）季度模型",
                        ["4-1-a 四季度交易量"],
                        [0])
    # 說明請後製 & 還缺要加入的df
    prs = addSlideDF(prs, 5, data_explain)
    # -------------------------------------
    # p6
    prs = addBulletPage(prs,"KDD4 交易模型（一）季度模型",
                        ["4-1-b 交易量圓餅圖"],
                        [0])
    # 說明請後製 
    prs.slides[6].shapes.add_picture("img/circle.png", Inches(2), Inches(2))
    # -------------------------------------
    # p7
    prs = addBulletPage(prs,"KDD4 交易模型（一）季度模型",
                        ["4-1-c 四季度平均單價"],
                        [0])
    # 說明請後製 & 還缺要加入的df
    prs = addSlideDF(prs, 7, data_explain)
    # -------------------------------------
    # p8
    prs = addBulletPage(prs,"KDD4 交易模型（一）季度模型",
                        ["4-1-d 平均單價折線圖"],
                        [0])
    # 說明請後製 
    prs.slides[8].shapes.add_picture("img/plot.png", Inches(2), Inches(2))
    # -------------------------------------
    # -------------------------------------
    # p9
    prs = addBulletPage(prs,"KDD4 交易模型（二）屋齡模型",
                        ["4-2-a 屋齡交易量"],
                        [0])
    AM = translate_raw.groupby("AR").agg({"地段": "nunique", "unit_price": "mean", "area":"count",
                                      }).reset_index()
    AM_col = ["屋齡範圍","地段数","房單價","交易量"]
    AM.columns = ["屋齡範圍","地段数","房單價","交易量"]
    # 說明請後製
    prs = add_personal_DF(prs,Inches(8.0),Inches(0.8), 9, AM,AM_col)
    # -------------------------------------
    # p10
    prs = addBulletPage(prs,"KDD4 交易模型（二）屋齡模型",
                        ["4-2-b 不同屋齡範圍的交易量"],
                        [0])
    # 說明請後製 
    prs.slides[10].shapes.add_picture("img/radio1.png", Inches(2), Inches(2))
    # -------------------------------------
    # p11
    prs = addBulletPage(prs,"KDD4 交易模型（二）屋齡模型",
                        ["4-2-c 不同屋齡範圍的房單價"],
                        [0])
    # 說明請後製 
    prs.slides[11].shapes.add_picture("img/radio2.png", Inches(2), Inches(2), Inches(2), Inches(2))
    # -------------------------------------
    # -------------------------------------
    #### p12~51  
    for i,place in enumerate (places):
        prs = addBulletPage(prs,"KDD4 （三）路段金額型",
                        ["4-3  "+place+"路段"],
                        [0])
        # 說明請後製 找幾個比較特殊的講講
        prs.slides[i+12].shapes.add_picture("img/bar/image"+str(i)+".png", Inches(3), Inches(2), Inches(7), Inches(5))
    # -------------------------------------
    # -------------------------------------
    # p 52
    prs = addBulletPage(prs,"KDD4 交易模型（四）單位金額",
                        ["4-4-a 平均每坪售價"],
                        [0])
    # 說明請後製  ****數學公式 還梅傑圖*****
    prs.slides[51].shapes.add_picture("img/math.png", Inches(4), Inches(4), Inches(3), Inches(2))
    # -------------------------------------
    # p 53
    prs = addBulletPage(prs,"KDD4 交易模型（四）單位金額",
                        ["4-4-b 用途 : 住家用"],
                        [0])
    a0,a1,a2,a3,a4 = detail(df,places,ages)
    p = aa(a0,ages)
    p = p.reset_index()
    ## get columns  
    p_col = np.array(p.columns)
    # 說明請後製
    prs = add_personal_DF(prs,Inches(8.0),Inches(0.8), 52, p,p_col)
    # -------------------------------------
    # p 54
    prs = addBulletPage(prs,"KDD4 交易模型（四）單位金額",
                        ["4-4-b 用途 : 商業用"],
                        [0])
    p = aa(a1,ages)
    p = p.reset_index()
    ## get columns   
    p_col = np.array(p.columns)
    # 說明請後製 
    prs = add_personal_DF(prs,Inches(8.0),Inches(0.8), 53, p,p_col)
    # -------------------------------------
    # p 55
    prs = addBulletPage(prs,"KDD4 交易模型（四）單位金額",
                        ["4-4-b 用途 : 辦公用"],
                        [0])
    p = aa(a2,ages)
    p = p.reset_index()
    ## get columns  
    p_col = np.array(p.columns)
    # 說明請後製
    prs = add_personal_DF(prs,Inches(8.0),Inches(0.8), 54, p,p_col)
    # -------------------------------------
    # p 56
    prs = addBulletPage(prs,"KDD4 交易模型（四）單位金額",
                        ["4-4-b 用途 : 住商用"],
                        [0])
    p = aa(a3,ages)
    p = p.reset_index()
    ## get columns  
    p_col = np.array(p.columns)
    # 說明請後製 
    prs = add_personal_DF(prs,Inches(8.0),Inches(0.8), 55, p,p_col)
    # -------------------------------------
    # p 57
    prs = addBulletPage(prs,"KDD4 交易模型（四）單位金額",
                        ["4-4-b 用途 : 工業用"],
                        [0])
    p = aa(a4,ages)
    p = p.reset_index()
    ## get columns    
    p_col = np.array(p.columns)
    # 說明請後製  & 還缺要加入的df
    prs = add_personal_DF(prs,Inches(8.0),Inches(0.8), 56, p,p_col)

    return prs



if __name__=="__main__":
    pass

    








